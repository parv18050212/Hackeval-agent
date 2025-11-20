import os
from typing import Dict, Any, List, Optional

try:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
except Exception:
    Presentation = None
    PP_PLACEHOLDER = None

try:
    import pypdf
except Exception:
    pypdf = None

try:
    import fitz  
except Exception:
    fitz = None

def _pct(n: int, d: int) -> float:
    return 0.0 if d <= 0 else (100.0 * n / d)

def _clip(v: float, lo: float, hi: float) -> float:
    return max(lo, min(hi, v))

class FormatAgent:
    """Heuristic format grader for decks. Output: {'score': 0..5, 'notes': [...]}"""
    def __init__(self):
        self.max_score = 5

    def _eval_pptx(self, file_path: str) -> Dict[str, Any]:
        notes: List[str] = []
        score = float(self.max_score)
        if Presentation is None:
            return {"score": 0, "notes": ["python-pptx not available"]}

        try:
            prs = Presentation(file_path)
        except Exception as e:
            return {"score": 0, "notes": [f"PPTX open error: {type(e).__name__}"]}

        n_slides = len(prs.slides)
        if n_slides == 0:
            return {"score": 0, "notes": ["Empty deck"]}

        title_count = 0
        fonts: List[str] = []
        title_sizes: List[float] = []
        body_sizes: List[float] = []
        colors: List[str] = []

        def _pt(sz) -> Optional[float]:
            try:
                return float(getattr(sz, "pt"))
            except Exception:
                return None

        for slide in prs.slides:
            slide_title_found = False
            for shp in slide.shapes:
                if getattr(shp, "has_text_frame", False):
                    tf = shp.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            fname = (getattr(getattr(run, "font", None), "name", None) or "").strip()
                            if fname:
                                fonts.append(fname)
                            fsize = _pt(getattr(getattr(run, "font", None), "size", None))
                            if fsize:
                                if para.level == 0 and (getattr(run.font, "bold", False) or fsize >= 24):
                                    title_sizes.append(fsize)
                                else:
                                    body_sizes.append(fsize)
                            col = getattr(getattr(run, "font", None), "color", None)
                            rgb = getattr(getattr(col, "rgb", None), "rgb", None) if col else None
                            if rgb:
                                colors.append(str(rgb))
                if getattr(shp, "is_placeholder", False) and PP_PLACEHOLDER:
                    try:
                        pht = shp.placeholder_format.type
                        if pht in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                            slide_title_found = True
                    except Exception:
                        pass
            if slide_title_found:
                title_count += 1

        if _pct(title_count, n_slides) < 70:
            score -= 1.0
            notes.append(f"Titles missing on many slides ({title_count}/{n_slides}).")

        if fonts:
            from collections import Counter
            fam_mode, fam_cnt = Counter(fonts).most_common(1)[0]
            if _pct(fam_cnt, len(fonts)) < 60:
                score -= 1.0
                notes.append("Multiple font families used inconsistently.")
            else:
                notes.append(f"Dominant font: {fam_mode}.")

        too_small_titles = sum(1 for v in title_sizes if v < 24)
        too_small_body = sum(1 for v in body_sizes if v < 16)
        if too_small_titles > 0:
            score -= 0.5
            notes.append("Some titles below 24pt.")
        if _pct(too_small_body, max(1, len(body_sizes))) > 30:
            score -= 0.5
            notes.append("Many body texts below 16pt.")

        if len(set(colors)) > 8:
            score -= 0.5
            notes.append("Too many distinct text colors; tighten the palette.")

        score = _clip(round(score, 1), 0.0, self.max_score)
        return {"score": int(round(score)), "notes": notes}

    def _eval_pdf(self, file_path: str, text_hint: str) -> Dict[str, Any]:
        notes: List[str] = []
        score = float(self.max_score)

        if pypdf:
            try:
                reader = pypdf.PdfReader(file_path)
                sizes = []
                for pg in reader.pages:
                    box = pg.mediabox
                    sizes.append((round(float(box.width), 1), round(float(box.height), 1)))
                if len(set(sizes)) > 1:
                    score -= 1.0
                    notes.append("Inconsistent page sizes.")
            except Exception as e:
                notes.append(f"PDF parse warning: {type(e).__name__}")
        else:
            notes.append("pypdf not available; limited checks.")

        text = text_hint or ""
        if not text and pypdf:
            try:
                reader = pypdf.PdfReader(file_path)
                text = "\n".join([(p.extract_text() or "") for p in reader.pages])
            except Exception:
                pass
        caps_lines = [ln for ln in (text or "").splitlines() if ln.strip() and ln.strip() == ln.strip().upper() and len(ln.strip()) >= 6]
        if len(caps_lines) < 3:
            score -= 0.5
            notes.append("Few obvious headings detected; add clear slide titles.")

        if fitz:
            try:
                doc = fitz.open(file_path)
                bg_means = []
                for page in doc[: min(len(doc), 6)]:
                    pix = page.get_pixmap(matrix=fitz.Matrix(1, 1))
                    import PIL.Image as PILImage
                    import io as _io
                    im = PILImage.open(_io.BytesIO(pix.tobytes("png"))).convert("L")
                    mean = sum(im.getdata()) / (im.width * im.height)
                    bg_means.append(round(mean, 1))
                if bg_means and (max(bg_means) - min(bg_means)) > 40:
                    score -= 0.5
                    notes.append("Large background variation; theme likely inconsistent.")
            except Exception:
                pass

        score = _clip(round(score, 1), 0.0, self.max_score)
        return {"score": int(round(score)), "notes": notes}

    async def aevaluate(self, file_path: str, text_hint: str = "") -> Dict[str, Any]:
        ext = os.path.splitext(file_path)[1].lower()
        if ext in (".pptx", ".ppt"):
            return self._eval_pptx(file_path)
        if ext == ".pdf":
            return self._eval_pdf(file_path, text_hint)
        return {"score": 0, "notes": ["Unsupported file type for format scoring."]}