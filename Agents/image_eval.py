# Robust image/diagram analysis without LangChain; uses OpenAI SDK directly.

import os
import io
import base64
from typing import List, Dict, Any, Optional

from pydantic import BaseModel, Field
from openai import OpenAI

# Prefer PyMuPDF for robust PDF image handling
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

# Fallback renderer
try:
    import pypdfium2 as pdfium
except Exception:
    pdfium = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from PIL import Image
except Exception:
    Image = None

from utils import extract_first_json_object as _extract_json_in_text, _render_pptx_slides_windows, _render_with_soffice

def _to_b64_jpeg_from_png_bytes(png_bytes: bytes, quality=85) -> str:
    if Image is None:
        return base64.b64encode(png_bytes).decode("utf-8")
    pil = Image.open(io.BytesIO(png_bytes))
    if pil.mode == "RGBA":
        pil = pil.convert("RGB")
    buf = io.BytesIO()
    pil.save(buf, format="JPEG", quality=quality)
    return base64.b64encode(buf.getvalue()).decode("utf-8")

class ImageAnalysis(BaseModel):
    image_index: int = Field(description="Index number of the image being analyzed.")
    description: str = Field(description="Step-by-step description of the diagram or image.")
    type: str = Field(description="Diagram type, e.g., Architecture, User Flow, Data Flow, Sequence, Chart, Mockup.")
    slide_index: Optional[int] = Field(default=None, description="0-based slide index if from PPT/PPTX.")
    page_index: Optional[int] = Field(default=None, description="0-based page index if from PDF.")
    is_diagram: bool = Field(description="True if the image is a diagram/flow/architecture, False if photo/logo.")
    importance: str = Field(description="One of: critical, supporting, decorative, irrelevant.")
    confidence: float = Field(description="0.0–1.0 confidence in the classification.", ge=0.0, le=1.0)

class WorkflowReport(BaseModel):
    overall_summary: str
    image_analyses: List[ImageAnalysis]

PROMPT_INSTRUCTIONS = """
You are a system design and process analysis specialist.
Analyze the images extracted from a presentation. These images can be flowcharts, architecture diagrams,
user journeys, data/ML pipelines, deployment diagrams, charts, mockups, or photos.

Instructions:
- If it's a diagram/flow/architecture, set is_diagram=true.
- importance: "critical" | "supporting" | "decorative" | "irrelevant".
- Keep descriptions concise and specific.
- Include slide/page indices when present.
- After all analyses, produce an overall workflow summary based on critical/supporting diagrams.

Return ONLY JSON matching this schema:

{
  "overall_summary": "string",
  "image_analyses": [
    {
      "image_index": 1,
      "description": "string",
      "type": "Architecture|User Flow|Data Flow|Sequence|Chart|Mockup|Photo|Other",
      "slide_index": 0,
      "page_index": 0,
      "is_diagram": true,
      "importance": "critical|supporting|decorative|irrelevant",
      "confidence": 0.7
    }
  ]
}
"""

class WorkflowAnalysisAgent:
    """
    Robust capture:
    - PDF: PyMuPDF page renders + embedded images as pixmaps to avoid PIL decode errors.
    - PPTX: embedded pictures + rendered slides via COM/soffice helpers.
    """

    def __init__(self):
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OPENAI_API_KEY not found in env")
        self.client = OpenAI(api_key=api_key)
        self.vision_model = os.getenv("OPENAI_MODEL_VISION", os.getenv("OPENAI_MODEL", "gpt-4o"))
        self.max_images = int(os.getenv("MAX_VISION_IMAGES", "12"))

    # -------- PDF helpers (PyMuPDF preferred) --------
    def _pdf_images_via_fitz(self, file_path: str) -> List[Dict[str, Any]]:
        if fitz is None:
            return []
        out: List[Dict[str, Any]] = []
        try:
            doc = fitz.open(file_path)
            max_pages = int(os.getenv("MAX_RENDER_PAGES", "12"))
            for pidx in range(min(len(doc), max_pages)):
                page = doc[pidx]
                # page render
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    out.append({"b64": _to_b64_jpeg_from_png_bytes(pix.tobytes("png")), "page_index": pidx})
                except Exception:
                    pass
                # embedded images
                try:
                    for img in page.get_images(full=True):
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.alpha:
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        out.append({"b64": _to_b64_jpeg_from_png_bytes(pix.tobytes("png")), "page_index": pidx})
                except Exception:
                    pass
        except Exception as e:
            print(f"[pdf fitz warn] {type(e).__name__}: {e}")
        return out

    def _pdf_pages_via_pdfium(self, file_path: str) -> List[Dict[str, Any]]:
        if pdfium is None:
            return []
        out: List[Dict[str, Any]] = []
        try:
            pdf = pdfium.PdfDocument(file_path)
            max_pages = int(os.getenv("MAX_RENDER_PAGES", "12"))
            for pidx in range(min(len(pdf), max_pages)):
                page = pdf[pidx]
                pil = page.render(scale=2.0).to_pil()
                if Image and pil.mode == "RGBA":
                    pil = pil.convert("RGB")
                buf = io.BytesIO()
                (pil or page.render(scale=2.0).to_pil()).save(buf, format="JPEG", quality=85)
                out.append({"b64": base64.b64encode(buf.getvalue()).decode("utf-8"), "page_index": pidx})
        except Exception as e:
            print(f"[pdf pdfium warn] {type(e).__name__}: {e}")
        return out

    # -------- PPTX helpers --------
    def _pptx_images(self, file_path: str) -> List[Dict[str, Any]]:
        out: List[Dict[str, Any]] = []
        if Presentation is None:
            return out
        try:
            prs = Presentation(file_path)
            for s_i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            if Image:
                                pil = Image.open(io.BytesIO(shape.image.blob))
                                if pil.mode == "RGBA":
                                    pil = pil.convert("RGB")
                                buf = io.BytesIO()
                                pil.save(buf, format="JPEG", quality=85)
                                b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
                            else:
                                b64 = base64.b64encode(shape.image.blob).decode("utf-8")
                            out.append({"b64": b64, "slide_index": s_i})
                        except Exception:
                            continue
        except Exception as e:
            print(f"[pptx warn] {type(e).__name__}: {e}")
        # also render slides via OS helpers if available
        try:
            rendered = _render_pptx_slides_windows(file_path) if os.name == "nt" else _render_with_soffice(file_path)
            for r in rendered[: int(os.getenv("MAX_RENDER_PAGES", "12"))]:
                out.append({"b64": r["b64"], "slide_index": r.get("slide_index")})
        except Exception:
            pass
        return out

    def _extract_images(self, file_path: str) -> List[Dict[str, Any]]:
        lower = file_path.lower()
        if lower.endswith(".pdf"):
            imgs = self._pdf_images_via_fitz(file_path)
            if not imgs:
                imgs = self._pdf_pages_via_pdfium(file_path)
            return imgs
        if lower.endswith(".pptx") or lower.endswith(".ppt"):
            return self._pptx_images(file_path)
        return []

    def analyze_workflows(self, file_path: str) -> Optional[WorkflowReport]:
        images = self._extract_images(file_path)
        if not images:
            print("  -> No images found to analyze.")
            return None

        # Build a single user message with alternating text + images
        # (OpenAI chat.completions supports image_url inputs)
        content: List[Dict[str, Any]] = [{"type": "text", "text": PROMPT_INSTRUCTIONS}]
        for idx, d in enumerate(images[: self.max_images], start=1):
            content.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{d['b64']}"}
            })
            ctx_bits = []
            if d.get("slide_index") is not None:
                ctx_bits.append(f"(slide {int(d['slide_index'])+1})")
            if d.get("page_index") is not None:
                ctx_bits.append(f"(page {int(d['page_index'])+1})")
            if ctx_bits:
                content.append({"type": "text", "text": f"Image {idx} context: {' '.join(ctx_bits)}"})

        try:
            resp = self.client.chat.completions.create(
                model=self.vision_model,
                messages=[{"role": "user", "content": content}],
                temperature=0.2,
            )
            raw = resp.choices[0].message.content or ""
            clean = _extract_json_in_text(raw) or raw
            import json
            data = json.loads(clean)
            # enrich indices & defaults
            enriched = []
            img_analyses = data.get("image_analyses", []) if isinstance(data, dict) else []
            for i, ia in enumerate(img_analyses, start=1):
                meta = images[i - 1] if i - 1 < len(images) else {}
                ia.setdefault("image_index", i)
                ia.setdefault("slide_index", meta.get("slide_index"))
                ia.setdefault("page_index", meta.get("page_index"))
                ia.setdefault("is_diagram", (ia.get("type", "").lower() not in ("photo", "image", "mockup", "logo")))
                ia.setdefault("importance", "supporting" if ia["is_diagram"] else "decorative")
                ia.setdefault("confidence", 0.7)
                enriched.append(ia)
            data["image_analyses"] = enriched
            return WorkflowReport.model_validate(data)
        except Exception as e:
            print(f"  -> ERROR during workflow analysis: {type(e).__name__}: {e}")
            return None