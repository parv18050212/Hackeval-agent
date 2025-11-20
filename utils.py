import os
import re
import json
import time
import shutil
import asyncio
import base64
import io
import tempfile
import glob as _glob
import subprocess
import hashlib
import math
from typing import List, Tuple, Dict, Any, Optional
from dataclasses import dataclass

from pptx import Presentation
import pypdf
from PIL import Image, ImageStat, UnidentifiedImageError

# Optional semantic search via OpenAI embeddings. If OPENAI_API_KEY is not set
# or embeddings calls fail, code falls back to simple keyword containment.
try:
    from openai import OpenAI
    _OPENAI_CLIENT = OpenAI(api_key=os.getenv("OPENAI_API_KEY")) if os.getenv("OPENAI_API_KEY") else None
    _EMBED_MODEL = os.getenv("OPENAI_EMBED_MODEL", "text-embedding-3-small")
except Exception:
    _OPENAI_CLIENT = None
    _EMBED_MODEL = None

# Optional imports used by the new evaluation-parameter loader
try:
    import openpyxl  # for Excel parsing
except Exception:
    openpyxl = None

try:
    import docx  # python-docx for DOCX parsing
except Exception:
    docx = None

try:
    import pypdfium2 as pdfium
except Exception:
    pdfium = None

try:
    import comtypes.client as _com_client
except Exception:
    _com_client = None

# ===================== Evaluation Schema =====================
@dataclass
class EvaluationParameters:
    criteria: Dict[str, Dict[str, Any]]
    rubric: str = ""

    @classmethod
    def default(cls):
        criteria = {
            "Problem identification & depth of understanding": {"weight": 10, "max_score": 10},
            "User-Centric Approach": {"weight": 10, "max_score": 10},
            "Innovation Quotient": {"weight": 15, "max_score": 10},
            "Technical Readiness & Prototype Potential": {"weight": 20, "max_score": 10},
            "Market Potential & Scalability": {"weight": 15, "max_score": 10},
            "Social/ Economic/ Environmental Impact": {"weight": 15, "max_score": 10},
            "Research depth & Ecosystem Awareness": {"weight": 10, "max_score": 10},
            "Presentation & Communication of Idea": {"weight": 5, "max_score": 10},
            "Format of the Presentation": {"weight": 5, "max_score": 10},
        }
        rubric = """
Use INTEGER 1–10 per criterion. Use full range. Avoid bunching.
Anchors:
10 exceptional with proof (metrics, full architecture or demo)
8 strong with one notable gap
6 adequate with multiple gaps; limited evidence
4 minimal coverage; mostly claims
2 not addressed

Rules:
- Penalize missing evidence. Default is 5 only if mixed evidence.
- Treat diagram evidence equal to text; if conflict, prefer diagram.
- At most one 10. Bias downward if uncertain.
- Do not give the same value to more than three criteria.
Checklist:
problem framing, baselines, datasets, KPIs & eval plan,
architecture & scalability, latency/cost, risks & mitigations, privacy/compliance,
security, deployment plan, adoption path.
"""
        return cls(criteria=criteria, rubric=rubric.strip())

    def to_dict(self):
        return {"criteria": self.criteria, "rubric": self.rubric}

    @classmethod
    def from_dict(cls, data: Dict[str, Any]):
        return cls(criteria=data.get("criteria", {}), rubric=data.get("rubric", ""))

ALLOWED_EXTS = {".pdf", ".pptx", ".ppt"}

def raw_total(scores: dict, eval_params: EvaluationParameters) -> float:
    return round(sum(float(scores.get(k, 0)) for k in eval_params.criteria.keys()), 1)

def weighted_total(scores: dict, eval_params: EvaluationParameters) -> float:
    total = 0.0
    for k, details in eval_params.criteria.items():
        v = scores.get(k, None)
        max_k = float(details.get("max_score", 10))
        w = float(details.get("weight", 0))
        x = 0.0 if v is None else max(0.0, min(max_k, float(v)))
        total += (x / max_k) * w
    return round(total, 2)

def _stable_small_jitter(team_name: str) -> float:
    h = hashlib.sha256(team_name.encode("utf-8")).hexdigest()
    val = int(h[:6], 16) / float(0xFFFFFF)
    return round(val * 0.009, 6)

def tie_break_key(scores: dict, team_name: str, eval_params: EvaluationParameters):
    def s(k):
        details = eval_params.criteria.get(k, {})
        max_k = float(details.get("max_score", 10))
        v = float(scores.get(k, 0) or 0.0)
        return max(0.0, min(1.0, v / max_k))
    
    primary_keys = list(eval_params.criteria.keys())
    
    score_list = [
        -weighted_total(scores, eval_params)
    ]
    for k in primary_keys:
        score_list.append(-s(k))
    
    score_list.extend([
        -_stable_small_jitter(team_name),
        team_name.lower(),
    ])
    return tuple(score_list)

# ---------- Rate limiters ----------
class _RateLimiter:
    def __init__(self, rpm: int):
        self.min_interval = 60.0 / max(1, rpm)
        self._last_ts = 0.0
        self._lock = asyncio.Lock()
    async def acquire(self):
        async with self._lock:
            now = time.perf_counter()
            wait = self._last_ts + self.min_interval - now
            if wait > 0:
                await asyncio.sleep(wait)
            self._last_ts = time.perf_counter()

_TEXT_LIMITER = _RateLimiter(int(os.getenv("RATE_LIMIT_RPM_TEXT", "18")))
_VISION_LIMITER = _RateLimiter(int(os.getenv("RATE_LIMIT_RPM_VISION", "6")))
def get_text_limiter(): return _TEXT_LIMITER
def get_vision_limiter(): return _VISION_LIMITER

# ---------- JSON extraction ----------
def extract_first_json_object(text: str) -> Optional[str]:
    if not text:
        return None
    if "```json" in text:
        try:
            return text.split("```json", 1)[1].split("```", 1)[0].strip()
        except Exception:
            pass
    start = -1
    depth = 0
    in_str = False
    esc = False
    for i, ch in enumerate(text):
        if in_str:
            if esc:
                esc = False
            elif ch == "\\":
                esc = True
            elif ch == '"':
                in_str = False
            continue
        if ch == '"':
            in_str = True
            continue
        if ch == "{":
            if depth == 0:
                start = i
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0 and start != -1:
                return text[start:i+1]
    return None

# ---------- Image helpers ----------
def _is_decorative(pil: Image.Image) -> bool:
    try:
        w, h = pil.size
        if w * h < 30000:
            return True
        stat = ImageStat.Stat(pil.convert("L"))
        if stat.var[0] < 50:
            return True
    except Exception:
        return False
    return False

def _to_b64_jpeg(pil_img: Image.Image, quality=85) -> str:
    if pil_img.mode == "RGBA":
        pil_img = pil_img.convert("RGB")
    buf = io.BytesIO()
    pil_img.save(buf, format="JPEG", quality=quality)
    return base64.b64encode(buf.getvalue()).decode("utf-8")

# ---------- PDF loaders ----------
def _render_pdf_pages_to_images(path: str) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    if pdfium is None:
        return out
    try:
        pdf = pdfium.PdfDocument(path)
        max_pages = int(os.getenv("MAX_RENDER_PAGES", "12"))
        for i in range(min(len(pdf), max_pages)):
            page = pdf[i]
            pil = page.render(scale=1.5).to_pil().convert("RGB")
            if not _is_decorative(pil):
                out.append({"b64": _to_b64_jpeg(pil), "page_index": i})
    except Exception as e:
        print(f"[pdf render warn] {type(e).__name__}")
    return out

def _extract_pdf_text_and_images(path: str) -> Tuple[str, List[str]]:
    text_parts: List[str] = []
    images_b64: List[str] = []
    with open(path, "rb") as f:
        reader = pypdf.PdfReader(f)
        for page in reader.pages:
            try:
                text_parts.append(page.extract_text() or "")
            except Exception:
                pass
        try:
            for page in reader.pages:
                if "/Resources" in page and "/XObject" in page["/Resources"]:
                    xobj = page["/Resources"]["/XObject"].get_object()
                    for name in xobj:
                        o = xobj[name]
                        if o.get("/Subtype") == "/Image":
                            try:
                                data = o.get_data()
                                pil = Image.open(io.BytesIO(data)).convert("RGB")
                                if not _is_decorative(pil):
                                    images_b64.append(_to_b64_jpeg(pil))
                            except UnidentifiedImageError:
                                continue
                            except Exception:
                                continue
        except Exception:
            pass
    pages = _render_pdf_pages_to_images(path)
    max_pages = int(os.getenv("MAX_RENDER_PAGES", "12"))
    if max_pages > 0:
        pages = pages[:max_pages]
    images_b64.extend([p["b64"] for p in pages])
    return "\n".join(text_parts), images_b64

# ---------- PPT/PPTX loaders ----------
def _render_pptx_slides_windows(path: str) -> List[Dict[str, Any]]:
    if _com_client is None or os.name != "nt":
        return []
    tmpdir = tempfile.mkdtemp(prefix="ppt_render_")
    out: List[Dict[str, Any]] = []
    pp = None
    pres = None
    try:
        pp = _com_client.CreateObject("PowerPoint.Application")
        pp.Visible = 0
        pres = pp.Presentations.Open(path, WithWindow=False)
        pres.Export(tmpdir, "PNG")
        files = sorted(_glob.glob(os.path.join(tmpdir, "*.PNG")) + _glob.glob(os.path.join(tmpdir, "*.png")))
        for idx, png in enumerate(files):
            try:
                pil = Image.open(png).convert("RGB")
                if not _is_decorative(pil):
                    out.append({"b64": _to_b64_jpeg(pil), "slide_index": idx})
            except Exception:
                pass
    except Exception as e:
        print(f"[pptx render error] {type(e).__name__}: {e}")
    finally:
        try:
            if pres: pres.Close()
        except Exception:
            pass
        try:
            if pp: pp.Quit()
        except Exception:
            pass
        shutil.rmtree(tmpdir, ignore_errors=True)
    return out

def _render_with_soffice(path: str) -> List[Dict[str, Any]]:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return []
    tmpdir = tempfile.mkdtemp(prefix="soffice_")
    out: List[Dict[str, Any]] = []
    try:
        cmd = [soffice, "--headless", "--convert-to", "png", "--outdir", tmpdir, path]
        subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
        files = sorted(_glob.glob(os.path.join(tmpdir, "*.png")))
        for idx, png in enumerate(files):
            try:
                pil = Image.open(png).convert("RGB")
                if not _is_decorative(pil):
                    out.append({"b64": _to_b64_jpeg(pil), "slide_index": idx})
            except Exception:
                pass
    except Exception as e:
        print(f"[soffice error] {type(e).__name__}: {e}")
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)
    return out

def _extract_pptx_text_and_images(path: str) -> Tuple[str, List[str]]:
    text_parts: List[str] = []
    images_b64: List[str] = []
    try:
        prs = Presentation(path)
    except Exception:
        prs = None
    if prs:
        for s_i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and getattr(shape, "has_text_frame", False):
                    try:
                        text_parts.append(shape.text)
                    except Exception:
                        pass
            for shape in slide.shapes:
                if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                    try:
                        pil = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                        if not _is_decorative(pil):
                            images_b64.append(_to_b64_jpeg(pil))
                    except Exception:
                        continue
    rendered = _render_pptx_slides_windows(path) if os.name == "nt" else _render_with_soffice(path)
    max_pages = int(os.getenv("MAX_RENDER_PAGES", "12"))
    if max_pages > 0:
        rendered = rendered[:max_pages]
    images_b64.extend([r["b64"] for r in rendered])
    return "\n".join(text_parts), images_b64

def load_document_content(file_path: str) -> Tuple[str, List[str]]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in ALLOWED_EXTS:
        return "", []
    if ext == ".pdf":
        return _extract_pdf_text_and_images(file_path)
    if ext in (".pptx", ".ppt"):
        return _extract_pptx_text_and_images(file_path)
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read(), []
    except Exception:
        return "", []

def _lower(s: Optional[str]) -> str:
    return (s or "").lower()

def _cosine(a: List[float], b: List[float]) -> float:
    try:
        da = math.sqrt(sum(x * x for x in a))
        db = math.sqrt(sum(x * x for x in b))
        if da == 0 or db == 0:
            return 0.0
        return sum(x * y for x, y in zip(a, b)) / (da * db)
    except Exception:
        return 0.0


def _hit_count(text: str, words: List[str]) -> int:
    """
    Prefer semantic matching (via OpenAI embeddings) when available. Returns
    how many of the provided word/phrase patterns are considered present in
    `text`.

    Fallback: simple substring containment for environments without API key
    or when embeddings calls fail.
    """
    if not text or not words:
        return 0

    # Try semantic route first if client is configured
    if _OPENAI_CLIENT and _EMBED_MODEL:
        try:
            # Build inputs: first the full text, then each phrase
            inputs = [text] + [str(w) for w in words]
            resp = _OPENAI_CLIENT.embeddings.create(model=_EMBED_MODEL, input=inputs)
            if not getattr(resp, "data", None):
                raise ValueError("no embedding data")
            emb_text = resp.data[0].embedding
            count = 0
            # threshold tuned to be permissive but avoid false positives
            threshold = float(os.getenv("SEMANTIC_MATCH_THRESHOLD", "0.72"))
            for i in range(1, len(inputs)):
                emb_w = resp.data[i].embedding
                sim = _cosine(emb_text, emb_w)
                if sim >= threshold:
                    count += 1
            return count
        except Exception:
            # If embeddings fail for any reason, fall through to keyword check
            pass

    # Fallback - original keyword containment (case-insensitive)
    t = _lower(text)
    return sum(1 for w in words if (w or "").lower() in t)

def _strength(count: int) -> int:
    if count <= 0: return 0
    if count == 1: return 1
    if count <= 3: return 2
    return 3

# ---------- Excel helpers ----------
def _excel_sanitize(value, limit=32000):
    def _clean_str(s: str) -> str:
        s = "".join(ch if (ch in "\r\n\t" or ord(ch) >= 32) else " " for ch in s)
        return s[:limit] if len(s) > limit else s

    if value is None:
        return None
    if isinstance(value, (int, float, bool)):
        return value
    if isinstance(value, str):
        return _clean_str(value)
    if isinstance(value, list):
        try:
            if all(isinstance(x, str) for x in value):
                return _clean_str("\n".join(x.strip() for x in value))
        except Exception:
            pass
        try:
            return _clean_str(json.dumps(value, ensure_ascii=False))
        except Exception:
            return _clean_str(str(value))
    if isinstance(value, dict):
        try:
            return _clean_str(json.dumps(value, ensure_ascii=False))
        except Exception:
            return _clean_str(str(value))
    try:
        return _clean_str(str(value))
    except Exception:
        return _clean_str(repr(value))

# ---------- Reporting ----------
def display_consolidated_report(context, eval_params: EvaluationParameters) -> None:
    print(f"\nTeam: {context.team_name}")
    if context.evaluation_error:
        print(f"Error: {context.evaluation_error}\n")
        return

    print("Scores:")
    for k in eval_params.criteria.keys():
        v = context.scores.get(k, "-")
        print(f"  - {k}: {v}")
    print(f"Total (raw): {raw_total(context.scores, eval_params)}")
    print(f"Total (weighted): {weighted_total(context.scores, eval_params)}")
    print(f"Summary: {context.scoring_summary or ''}")

    overall = None
    if isinstance(context.workflow_analysis, dict):
        overall = context.workflow_analysis.get("overall_summary")
    if not overall and getattr(context, "workflow_report", None):
        overall = context.workflow_report.get("overall_summary")
    if not overall:
        overall = "Null. No diagrams were provided in the presentation text."
    print(f"Workflow summary: {overall}")

    if context.feedback:
        print("Feedback:")
        print(f"  Positive: {context.feedback.get('positive','')}")
        print(f"  Criticism: {context.feedback.get('criticism','')}")
        print(f"  Technical: {context.feedback.get('technical','')}")
        print(f"  Suggestions: {context.feedback.get('suggestions','')}")
    print()

def display_leaderboard(contexts: list, eval_params: EvaluationParameters) -> None:
    printable = []
    for c in contexts:
        total = weighted_total(c.scores, eval_params) if not c.evaluation_error else -1
        printable.append((total, c.team_name, c.file_path, c.scores or {}))
    printable.sort(key=lambda x: tie_break_key(x[3], x[1], eval_params))
    
    print("\n######## Leaderboard ########")
    for i, (total, name, path, scores) in enumerate(printable, 1):
        status = f"{total:.2f}" if total >= 0 else "ERROR"
        score_details = " | ".join([
            f"{k.split()[0]} {scores.get(k, '-')}" for k in eval_params.criteria.keys()[:3]
        ])
        print(f"{i:2d}. {name:25s} | weighted {status:>6s} | {score_details} | {os.path.basename(path)}")
    print("#############################\n")

def save_consolidated_reports_to_excel(contexts: list, filename: str, eval_params: EvaluationParameters):
    try:
        from openpyxl import Workbook
    except ImportError:
        print("openpyxl not installed. Skipping Excel export.")
        return

    wb = Workbook()
    ws = wb.active
    ws.title = "Reports"
    headers = [
        "team_name","file_path","evaluation_error"
    ] + list(eval_params.criteria.keys()) + [
        "total_raw","total_weighted","summary","workflow_overall",
        "feedback_positive","feedback_criticism","feedback_technical","feedback_suggestions"
    ]
    ws.append(headers)

    for ctx in contexts:
        row = [
            _excel_sanitize(getattr(ctx, "team_name", "")),
            _excel_sanitize(getattr(ctx, "file_path", "")),
            _excel_sanitize(getattr(ctx, "evaluation_error", "")),
        ]

        for k in eval_params.criteria.keys():
            row.append(_excel_sanitize(ctx.scores.get(k, "-") if ctx.scores else "-"))

        row.append(_excel_sanitize(raw_total(ctx.scores, eval_params) if ctx.scores else "-"))
        row.append(_excel_sanitize(weighted_total(ctx.scores, eval_params) if ctx.scores else "-"))

        row.append(_excel_sanitize(getattr(ctx, "scoring_summary", "")))

        overall = None
        if isinstance(ctx.workflow_analysis, dict):
            overall = ctx.workflow_analysis.get("overall_summary")
        if not overall and getattr(ctx, "workflow_report", None):
            overall = ctx.workflow_report.get("overall_summary")
        if not overall:
            overall = "Null. No diagrams were provided in the presentation text."
        row.append(_excel_sanitize(overall))

        fb = ctx.feedback or {}
        row += [
            _excel_sanitize(fb.get("positive", "")),
            _excel_sanitize(fb.get("criticism", "")),
            _excel_sanitize(fb.get("technical", "")),
            _excel_sanitize(fb.get("suggestions", "")),
        ]

        ws.append(row)

    wb.save(filename)
    print(f"[Excel] Consolidated reports saved to {filename}")

def save_leaderboard_to_excel(contexts: list, filename: str, eval_params: EvaluationParameters):
    try:
        from openpyxl import Workbook
    except ImportError:
        print("openpyxl not installed. Skipping Excel export.")
        return

    wb = Workbook()
    ws = wb.active
    ws.title = "Leaderboard"
    headers = ["Rank", "Team Name", "Weighted Total"] + list(eval_params.criteria.keys()) + ["File Name"]
    ws.append(headers)

    printable = []
    for c in contexts:
        total = weighted_total(c.scores, eval_params) if not c.evaluation_error else -1
        printable.append((total, c.team_name, c.file_path, c.scores or {}))
    printable.sort(key=lambda x: tie_break_key(x[3], x[1], eval_params))

    for i, (total, name, path, scores) in enumerate(printable, 1):
        row = [
            i,
            _excel_sanitize(name),
            _excel_sanitize(f"{total:.2f}" if total >= 0 else "ERROR"),
        ]
        for k in eval_params.criteria.keys():
            row.append(_excel_sanitize(scores.get(k, "-")))
        row.append(_excel_sanitize(os.path.basename(path)))
        ws.append(row)

    wb.save(filename)
    print(f"[Excel] Leaderboard saved to {filename}")

# ===================== New: Evaluation Parameters Loader (JSON/PDF/Excel/DOCX) =====================

def _normalize_criteria_dict(raw: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
    """
    Normalize incoming criteria dict to {name: {weight: int, max_score: int}}.
    Accepts flexible key casing and synonyms like 'max' for 'max_score'.
    """
    out: Dict[str, Dict[str, Any]] = {}
    for k, v in (raw or {}).items():
        if not isinstance(v, dict):
            # allow short format like {"Criterion": 10} => interpret as max_score with default weight 10
            try:
                num = int(v)
                out[k] = {"weight": 10, "max_score": num}
                continue
            except Exception:
                continue
        weight = v.get("weight", v.get("Weight", v.get("w", 10)))
        max_score = v.get("max_score", v.get("max", v.get("Max", 10)))
        try:
            weight = int(weight)
        except Exception:
            weight = 10
        try:
            max_score = int(max_score)
        except Exception:
            max_score = 10
        out[k] = {"weight": weight, "max_score": max_score}
    return out

def _parse_eval_params_from_json_text(text: str) -> Optional[EvaluationParameters]:
    try:
        obj = json.loads(text)
        criteria = _normalize_criteria_dict(obj.get("criteria", {}))
        rubric = (obj.get("rubric") or "").strip()
        if criteria:
            return EvaluationParameters(criteria=criteria, rubric=rubric)
    except Exception:
        pass
    # If the whole object is criteria (no top-level), accept that too
    try:
        obj = json.loads(text)
        if isinstance(obj, dict):
            criteria = _normalize_criteria_dict(obj)
            if criteria:
                return EvaluationParameters(criteria=criteria, rubric="")
    except Exception:
        pass
    return None

def _parse_eval_params_from_freeform_text(text: str) -> Optional[EvaluationParameters]:
    """
    Strategy:
    1) Try to extract an embedded JSON object and parse.
    2) Else, parse lines heuristically for 'criterion | weight | max' style.
    3) Else, return None.
    """
    if not text:
        return None

    # 1) Embedded JSON
    js = extract_first_json_object(text)
    if js:
        parsed = _parse_eval_params_from_json_text(js)
        if parsed:
            return parsed

    # 2) Heuristic line parser: look for rows like
    # "Criterion Name | weight=15 | max=10" or "Criterion,15,10" or "Criterion - weight:15, max:10"
    criteria: Dict[str, Dict[str, Any]] = {}
    rubric_lines: List[str] = []
    for ln in text.splitlines():
        raw = ln.strip()
        if not raw:
            continue

        # Collect rubric following an explicit "Rubric" header
        if re.match(r"^\s*rubric\s*[:\-]", raw, flags=re.I):
            rub = raw.split(":", 1)[-1] if ":" in raw else raw.split("-", 1)[-1]
            rubric_lines.append(rub.strip())
            continue

        # Attempt various formats
        m = re.match(r"^\s*[-*]?\s*(.+?)\s*[\|\-,:]\s*(?:weight\s*[:=]\s*(\d+))\s*[,| ]+\s*(?:max(?:_score)?\s*[:=]\s*(\d+))\s*$", raw, flags=re.I)
        if not m:
            # CSV-ish: Criterion,15,10
            m = re.match(r"^\s*[-*]?\s*(.+?)\s*[,|]\s*(\d+)\s*[,|]\s*(\d+)\s*$", raw)
        if m:
            name = m.group(1).strip()
            w = int(m.group(2))
            mx = int(m.group(3))
            if name:
                criteria[name] = {"weight": w, "max_score": mx}
            continue

    if criteria:
        rubric = "\n".join(rubric_lines).strip()
        return EvaluationParameters(criteria=criteria, rubric=rubric)

    return None

def _load_eval_params_from_json_file(path: str) -> Optional[EvaluationParameters]:
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        criteria = _normalize_criteria_dict(data.get("criteria", data if isinstance(data, dict) else {}))
        rubric = (data.get("rubric") or "") if isinstance(data, dict) else ""
        if criteria:
            return EvaluationParameters(criteria=criteria, rubric=(rubric or "").strip())
    except Exception as e:
        print(f"[eval-params] JSON parse error: {type(e).__name__}: {e}")
    return None

def _load_eval_params_from_pdf(path: str) -> Optional[EvaluationParameters]:
    try:
        with open(path, "rb") as f:
            reader = pypdf.PdfReader(f)
            text = "\n".join([(p.extract_text() or "") for p in reader.pages])
        return _parse_eval_params_from_freeform_text(text)
    except Exception as e:
        print(f"[eval-params] PDF parse error: {type(e).__name__}: {e}")
        return None

def _load_eval_params_from_docx(path: str) -> Optional[EvaluationParameters]:
    if docx is None:
        print("[eval-params] python-docx not installed; cannot parse DOCX.")
        return None
    try:
        d = docx.Document(path)
        parts: List[str] = []
        # Paragraph text
        for p in d.paragraphs:
            parts.append(p.text or "")
        # Table text
        for t in d.tables:
            for row in t.rows:
                cells = [c.text.strip() for c in row.cells]
                parts.append(" | ".join(cells))
        text = "\n".join(parts)
        parsed = _parse_eval_params_from_freeform_text(text)
        if parsed:
            return parsed
    except Exception as e:
        print(f"[eval-params] DOCX parse error: {type(e).__name__}: {e}")
    return None

def _load_eval_params_from_excel(path: str) -> Optional[EvaluationParameters]:
    if openpyxl is None:
        print("[eval-params] openpyxl not installed; cannot parse Excel.")
        return None
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        # Heuristic: first, try a sheet explicitly named like 'criteria' or 'evaluation'
        candidate_sheets = []
        for ws in wb.worksheets:
            nm = (ws.title or "").lower()
            if any(key in nm for key in ["criteria", "evaluation", "scoring", "rubric"]):
                candidate_sheets.append(ws)
        if not candidate_sheets:
            candidate_sheets = [wb.active]

        def _read_sheet(ws) -> Tuple[Dict[str, Dict[str, Any]], str]:
            criteria: Dict[str, Dict[str, Any]] = {}
            rubric = ""
            # Find a header row
            header_row_idx = None
            headers: Dict[int, str] = {}
            for r in range(1, min(ws.max_row, 25) + 1):
                row_vals = [str(ws.cell(r, c).value).strip() if ws.cell(r, c).value is not None else "" for c in range(1, min(ws.max_column, 10) + 1)]
                lower = [v.lower() for v in row_vals]
                if any("criterion" in v or "criteria" in v for v in lower) and any("weight" in v for v in lower) and any("max" in v for v in lower):
                    header_row_idx = r
                    for c, v in enumerate(row_vals, start=1):
                        headers[c] = v.strip().lower()
                    break
            # If header found, parse rows
            if header_row_idx:
                for r in range(header_row_idx + 1, ws.max_row + 1):
                    row = {headers.get(c, ""): ws.cell(r, c).value for c in headers.keys()}
                    name = str(row.get("criterion") or row.get("criteria") or row.get("name") or "").strip()
                    if not name:
                        continue
                    try:
                        w = int(float(row.get("weight") or 10))
                    except Exception:
                        w = 10
                    try:
                        mx = int(float(row.get("max_score") or row.get("max") or 10))
                    except Exception:
                        mx = 10
                    criteria[name] = {"weight": w, "max_score": mx}
            # Try to get rubric either from a 'rubric' column or a dedicated sheet
            # 1) single cell labeled 'rubric'
            for r in range(1, min(ws.max_row, 50) + 1):
                for c in range(1, min(ws.max_column, 10) + 1):
                    val = ws.cell(r, c).value
                    if isinstance(val, str) and val.strip().lower().startswith("rubric"):
                        # next cell to the right if exists
                        nxt = ws.cell(r, c + 1).value
                        if isinstance(nxt, str):
                            rubric = nxt.strip()
                        else:
                            rubric = val.split(":", 1)[-1].strip()
                        break
            return criteria, rubric

        all_criteria: Dict[str, Dict[str, Any]] = {}
        rubric_collected = ""
        for ws in candidate_sheets:
            c, r = _read_sheet(ws)
            all_criteria.update(c)
            if r and not rubric_collected:
                rubric_collected = r

        if all_criteria:
            return EvaluationParameters(criteria=all_criteria, rubric=(rubric_collected or "").strip())
    except Exception as e:
        print(f"[eval-params] Excel parse error: {type(e).__name__}: {e}")
    return None

def load_eval_parameters_from_path(path: Optional[str]) -> EvaluationParameters:
    """
    Public entry: load evaluation parameters from JSON, PDF, Excel, DOCX, or freeform text
    (if someone points to a .txt). Falls back to defaults on failure.
    """
    if not path or not os.path.exists(path):
        print("No evaluation parameters file found. Using default parameters.")
        return EvaluationParameters.default()

    ext = os.path.splitext(path)[1].lower()
    loaders = []
    if ext == ".json":
        loaders = [_load_eval_params_from_json_file]
    elif ext in (".pdf",):
        loaders = [_load_eval_params_from_pdf]
    elif ext in (".xlsx", ".xls"):
        loaders = [_load_eval_params_from_excel]
    elif ext in (".docx",):
        loaders = [_load_eval_params_from_docx]
    else:
        # Attempt freeform text, then JSON
        loaders = []

    for fn in loaders:
        res = fn(path)
        if res:
            return res

    # Fallback: try reading as text then parse
    try:
        with open(path, "r", encoding="utf-8") as f:
            txt = f.read()
        res = _parse_eval_params_from_freeform_text(txt) or _parse_eval_params_from_json_text(txt)
        if res:
            return res
    except Exception:
        pass

    print("Could not parse evaluation parameters from path. Using default parameters.")
    return EvaluationParameters.default()
